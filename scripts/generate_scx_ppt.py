#!/usr/bin/env python
"""Generate SCX academic presentation PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──
W = Inches(13.333)  # 16:9
H = Inches(7.5)
BLUE = RGBColor(0x1A, 0x56, 0xDB)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x0D, 0x94, 0x8B)  # teal accent
RED_ACCENT = RGBColor(0xDC, 0x26, 0x26)
ORANGE = RGBColor(0xE6, 0x8A, 0x00)
LIGHT_BLUE = RGBColor(0xDB, 0xE9, 0xFA)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── Helpers ──
def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_bg(slide, color=LIGHT_GRAY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multiline(slide, left, top, width, height, lines, font_size=16, color=DARK, line_spacing=1.5, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
    return tf

def add_rect(slide, left, top, width, height, fill_color=WHITE, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

def add_thin_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    return shape

# ═══════════════════════════════════════════════
# SLIDE 1: Cover
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_bg(slide, WHITE)

# Top accent bar
add_thin_line(slide, 0, 0, 13.333)

# Main title
add_textbox(slide, 1.5, 1.8, 10.3, 1.2,
    "SCX：状态条件数据价值评估框架", font_size=40, bold=True, color=DARK)

# Subtitle
add_textbox(slide, 1.5, 3.0, 10.3, 0.8,
    "The value of data is state-conditioned", font_size=24, color=BLUE, bold=False)

# Divider
add_thin_line(slide, 1.5, 3.9, 3.0)

# Description
add_textbox(slide, 1.5, 4.2, 10.3, 0.6,
    "从机器学习势函数到科学与工程仿真", font_size=20, color=GRAY)

# Author info
add_multiline(slide, 1.5, 5.5, 5, 1.0, [
    "汇报人：许少杰",
    "单位：华中科技大学",
    "日期：2026年6月",
], font_size=16, color=GRAY)

add_notes(slide, """各位老师好，我今天汇报的题目是《SCX：状态条件数据价值评估框架》。

这个工作的出发点很简单：在科学计算和机器学习里，数据很贵，但我们经常不知道哪些数据值得花钱去算、哪些是重复的、哪些可能有问题。

我最早是在做机器学习势函数的时候发现了这个问题的一些规律，后来发现这个问题结构在很多领域都存在。所以我把这个想法抽象成了一个通用框架，取名叫 SCX。

今天我主要讲四件事：问题是什么、核心想法是什么、已经做了什么、下一步想做什么。""")

# ═══════════════════════════════════════════════
# SLIDE 2: Problem Background
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_bg(slide, WHITE)
add_thin_line(slide, 0, 0, 13.333)

add_textbox(slide, 0.8, 0.5, 11.7, 0.7, "问题背景：数据很贵，但我们经常不知道哪些数据值得用", font_size=28, bold=True, color=DARK)

# Four problem cards
problems = [
    ("💰 数据采集成本高", "DFT一个结构可能算几十核时\nFEM一次仿真跑几小时\n医学图像标注需要专家"),
    ("❓ 数据价值不明确", "哪些数据值得继续算？\n哪些是重复的、可以跳过的？\n哪些数据可能是错的？"),
    ("🔀 多专家/多模型不可比", "不同专家在不同场景表现差异大\n没有一个全局最优的专家\n怎么判断当前该信哪个模型？"),
    ("📊 传统方法有局限", "LOO/Shapley给每个样本一个固定分数\n不考虑状态上下文\n不区分噪声和困难样本"),
]

for i, (title, body) in enumerate(problems):
    x = 0.5 + i * 3.1
    card = add_rect(slide, x, 1.5, 2.9, 3.8, fill_color=LIGHT_GRAY)
    add_textbox(slide, x + 0.2, 1.65, 2.5, 0.5, title, font_size=16, bold=True, color=DARK)
    add_multiline(slide, x + 0.2, 2.2, 2.5, 3.0, body.split('\n'), font_size=13, color=GRAY)

# Bottom summary
add_textbox(slide, 0.8, 5.8, 11.7, 1.0,
    "核心矛盾：我们需要一个框架，能在\"状态\"的粒度上判断——这个数据值不值得用、这个专家值不值得信、这个动作值不值得花钱。",
    font_size=18, bold=False, color=DARK)

add_notes(slide, """先说一下为什么要做这个。

在科学计算里，数据是很贵的。比如 DFT 算一个 AlN 结构要好几十核时，FEM 仿真一次跑几个小时，医学图像标注需要放射科医生一个一个看。但问题是——我们花大价钱算出来的数据，到底哪些是真的有用、哪些是重复劳动、哪些甚至是有问题的？

传统方法，比如 LOO、Data Shapley，它们给每个样本算一个固定分数，说"这个样本价值0.8，那个0.3"。但这个思路有个根本问题——数据价值不是样本的固有属性。同一个样本，在不同的模型、不同的已有数据覆盖下，价值完全不一样。

打个比方：一道物理题，对一个已经刷了100道同类题的学生来说是冗余的，但对一个刚入门的学生来说就是高价值的。数据价值取决于"谁在用、已经有什么、当前缺什么"。

这就是 SCX 要解决的核心矛盾。""")

# ═══════════════════════════════════════════════
# SLIDE 3: Core Idea
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_bg(slide, WHITE)
add_thin_line(slide, 0, 0, 13.333)

add_textbox(slide, 0.8, 0.5, 11.7, 0.7, "核心思想：数据价值是状态条件的", font_size=28, bold=True, color=DARK)

# Left: traditional
add_rect(slide, 0.5, 1.5, 5.8, 2.2, fill_color=LIGHT_GRAY)
add_textbox(slide, 0.8, 1.6, 5.2, 0.5, "传统方法：全局打分", font_size=20, bold=True, color=RED_ACCENT)
add_multiline(slide, 0.8, 2.1, 5.2, 1.3, [
    "• LOO / Data Shapley / Influence Function",
    "• 每个样本 → 一个固定分数",
    "• 不考虑：谁在用、已有啥、当前缺啥",
    "• 结果：噪声和困难样本分不清",
], font_size=14, color=GRAY)

# Right: SCX
add_rect(slide, 6.8, 1.5, 5.8, 2.2, fill_color=LIGHT_BLUE)
add_textbox(slide, 7.1, 1.6, 5.2, 0.5, "SCX：状态条件打分", font_size=20, bold=True, color=BLUE)
add_multiline(slide, 7.1, 2.1, 5.2, 1.3, [
    "• V(x) = V(x | state, coverage, expert, cost, feedback)",
    "• 同一个样本 → 不同条件下价值不同",
    "• 数据价值由五个因素共同决定",
    "• 可以区分：冗余 / 噪声 / 高价值 / 需路由",
], font_size=14, color=DARK)

# Formula
add_rect(slide, 0.5, 4.1, 12.3, 1.3, fill_color=WHITE, border_color=BLUE)
add_textbox(slide, 0.8, 4.2, 11.7, 0.5, "核心公式", font_size=16, bold=True, color=BLUE)
add_textbox(slide, 0.8, 4.7, 11.7, 0.6,
    "V(s) = r̄(s) · ρ(s) · L(s) · [1 − D(s)] · max_m SCX_m(s)", font_size=22, bold=True, color=DARK)
add_multiline(slide, 0.8, 5.1, 11.7, 0.5, [
    "r̄=当前误差   ρ=状态出现概率   L=可学习性(非噪声)   D=已有数据覆盖度   SCX=最佳专家在该状态的可靠性",
], font_size=11, color=GRAY)

# Key insight
add_textbox(slide, 0.8, 5.9, 11.7, 1.0,
    "关键洞察：数据价值不是样本的固有属性。同一批数据，在模型 A 眼里是噪声，在模型 B 手里可能是高价值样本。",
    font_size=17, bold=False, color=DARK)

add_notes(slide, """这一页是核心思想。

传统方法给每个样本一个固定的分数，比如 Data Shapley 算出来样本 A 价值 0.8，样本 B 价值 0.3。这个数字是"绝对"的，跟谁在用、已经有什么数据、当前缺什么都无关。

SCX 的核心思想很简单：数据价值取决于状态。这个"状态"包括很多因素——当前模型在这个区域的误差有多大、这个区域出现了多少次、里面的标签是不是一致的、已经有多少类似数据了、哪个专家在这个区域最可靠。

写成公式就是下面这个：V(s) 等于五个因子的乘积。任何一个因子为零，整个价值就为零。比如：如果当前误差已经很低了(r̄→0)，那这个状态的数据就不值钱——模型已经学会了。如果噪声很高(L→0)，也不值钱——再多的坏数据也学不出好东西。

这个思想的来源其实很朴素：在做势函数的时候发现，同一个模型，在声子附近预测得很准，在高温热振动附近预测得很差。那显然，你在声子区域继续加数据是浪费，而高温区域的数据质量参差不齐，有些是真实物理、有些就是脚本生成的坏构型。

SCX 把这个经验观察形式化了。""")

# ═══════════════════════════════════════════════
# SLIDE 4: Four-quadrant explanation with AlN data
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_bg(slide, WHITE)
add_thin_line(slide, 0, 0, 13.333)

add_textbox(slide, 0.8, 0.5, 11.7, 0.7, "大白话解释：同一个样本，在不同状态下价值完全不同", font_size=26, bold=True, color=DARK)

# Four-quadrant concept
quads = [
    ("冗余数据 (Redundant)", "低误差 + 高密度", "→ 模型已经学会了，可以压缩", "Phonon 120帧: fmax~0.66\n预估可压缩 50-80%", LIGHT_GRAY),
    ("高价值数据 (Valuable)", "高误差 + 高密度 + 高一致性", "→ 值得花钱补充训练", "EOS极端应变区\n弹性变形区", LIGHT_BLUE),
    ("噪声/坏标签 (Noisy)", "高误差 + 低密度 + 低一致性", "→ 需要重算或降权", "Thermal 1800K: 49.1%帧 fmax>5\n14帧 fmax>10 (极端异常)", RGBColor(0xFD, 0xE8, 0xE8)),
    ("需路由数据 (Expert-dep.)", "某专家在此状态显著优于其他", "→ 交给最合适的专家处理", "MLMD stress10 区域\n需路由到高温训练专家", RGBColor(0xFE, 0xF3, 0xE6)),
]

for i, (title, cond, action, example, bg) in enumerate(quads):
    x = 0.3 + (i % 2) * 6.5
    y = 1.4 + (i // 2) * 2.8
    card = add_rect(slide, x, y, 6.2, 2.5, fill_color=bg)
    add_textbox(slide, x + 0.2, y + 0.1, 5.8, 0.4, title, font_size=18, bold=True, color=DARK)
    add_textbox(slide, x + 0.2, y + 0.55, 5.8, 0.3, cond, font_size=13, color=BLUE)
    add_textbox(slide, x + 0.2, y + 0.9, 5.8, 0.3, action, font_size=14, bold=True, color=DARK)
    add_multiline(slide, x + 0.2, y + 1.3, 5.8, 1.2, example.split('\n'), font_size=12, color=GRAY)

add_notes(slide, """这一页用大白话解释 SCX 的四分类逻辑。

先说冗余数据：模型已经学得很好了，误差很低，而且类似的样本你已经有很多了。这种情况下再加数据就是浪费钱。在 AlN 的 phonon batch 里，120 帧全部 fmax 在 0.66 左右，模型对声子附近的力预测已经非常准。SCX 建议这个区域可以压缩 50-80%。

再说高价值数据：误差还比较高，但数据的一致性很好——也就是说，这些不是噪声，是真的还没学会。比如 EOS 的极端应变区、弹性变形的某些方向。这些地方值得继续投入。

然后是最关键的噪声数据：误差很高，但密度很低、一致性也低。这种大概率是数据本身有问题，不是模型的问题。在我们的 AlN 计算里，1800K 的热振动 snapshot 中有将近一半的帧 fmax 超过了 5 eV/Å，甚至有 14 帧超过了 10——物理上，一个 AlN 晶体里的原子受力不可能这么大。这就是 VASP 脚本生成的构型有问题，属于"数据中毒"。

最后是专家依赖的数据：某类数据，专家 A 处理得很好，专家 B 处理得很差。这时候不应该让两个专家平均投票，而应该把这类数据路由给 A。

这四个分类不是人为划的，是 SCX 根据状态条件自动判定的。""")

# ═══════════════════════════════════════════════
# SLIDE 5: SCX Four-classification
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_bg(slide, WHITE)
add_thin_line(slide, 0, 0, 13.333)

add_textbox(slide, 0.8, 0.5, 11.7, 0.7, "SCX 把数据分成四类，每类对应不同动作", font_size=28, bold=True, color=DARK)

# Table visualization
classes = [
    ("Valuable\n有价值的", "高误差 + 高密度\n+ 高一致性 + 低覆盖", "→ 采集更多", "EOS极端应变、弹性变形\n未覆盖的构型区域", BLUE),
    ("Redundant\n冗余的", "低误差 + 高密度\n+ 高覆盖", "→ 压缩/跳过", "Phonon小位移 120帧\n压缩 50-80% 精度不降", ACCENT),
    ("Noisy\n噪声风险的", "高误差 + 低密度\n+ 低一致性", "→ 重算/降权/丢弃", "Thermal 1800K 49%异常\nMLMD stress10 35%异常", RED_ACCENT),
    ("Expert-dep.\n专家依赖的", "某专家在此状态\n显著优于其他", "→ 路由到该专家", "MLMD 高温区需专项专家\n表面/缺陷需迁移学习", ORANGE),
]

for i, (name, cond, action, example, accent_color) in enumerate(classes):
    y = 1.4 + i * 1.45
    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.08), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Name
    add_textbox(slide, 0.8, y + 0.05, 1.8, 1.1, name, font_size=18, bold=True, color=DARK)
    # Condition
    add_textbox(slide, 2.8, y + 0.05, 3.2, 1.1, cond, font_size=14, color=GRAY)
    # Action
    add_textbox(slide, 6.2, y + 0.05, 2.5, 1.1, action, font_size=16, bold=True, color=accent_color)
    # Example
    add_textbox(slide, 9.0, y + 0.05, 4.0, 1.1, example, font_size=13, color=GRAY)

add_notes(slide, """这一页是四分类的具体标准和动作建议。

每个分类的判断条件已经很清楚了。我想强调的是"动作"这一列：

对于有价值的，动作是"采集更多"——这类数据是模型还没学会的，而且数据质量可靠，你应该在这个区域继续投入计算资源。

对于冗余的，动作是"压缩"——模型已经学会了，你再算一百个类似的结构也是浪费钱。SCX 在 AlN 上发现 phonon 的 120 个帧可以压缩掉一半以上。

对于噪声风险的，动作是"重算或降权"——这些数据本身有问题，强行用来训练反而会损害模型。SCX 能自动发现这些帧，你不需要一个一个去检查 fmax。

对于专家依赖的，动作是"路由"——别搞平均主义。在状态 A 专家 1 最强，就把状态 A 的数据交给专家 1。这对多专家系统、多保真度仿真系统都是直接可用的。""")

# ═══════════════════════════════════════════════
# SLIDE 6: Why MLIP first
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_bg(slide, WHITE)
add_thin_line(slide, 0, 0, 13.333)

add_textbox(slide, 0.8, 0.5, 11.7, 0.7, "为什么 MLIP 是 SCX 的第一个根据地", font_size=28, bold=True, color=DARK)

reasons = [
    ("❶ 天然的状态编码器", "ACE descriptor 把任意原子构型映射到 182 维基函数空间。\n这个映射本身就是一种状态编码——构型相似 → 基函数接近。\n不需要额外训练 embedding 网络。"),
    ("❷ 可靠的参考标签", "DFT 虽然贵，但提供了相对可靠的能量/力/应力标签。\n有 ground truth，才能画出 residual map，才能判断专家好不好。\n这比纯无监督场景多了一个锚点。"),
    ("❸ 天然的多专家场景", "同一个训练集，可以训练多个不同架构/不同初始化/不同数据子集的势函数。\n它们在同一个构型区域上的表现天然不同——这正是 SCX 需要对比的。"),
]

for i, (title, body) in enumerate(reasons):
    y = 1.4 + i * 1.7
    add_rect(slide, 0.5, y, 12.3, 1.5, fill_color=LIGHT_GRAY if i % 2 == 0 else LIGHT_BLUE)
    add_textbox(slide, 0.8, y + 0.1, 11.7, 0.4, title, font_size=20, bold=True, color=DARK)
    add_multiline(slide, 0.8, y + 0.55, 11.7, 0.9, body.split('\n'), font_size=14, color=GRAY)

# Bottom summary
add_textbox(slide, 0.8, 6.5, 11.7, 0.5,
    "总结：MLIP = 天然编码器(ACE) + 可靠标签(DFT) + 多专家对比 → SCX 的理想起点",
    font_size=17, bold=True, color=BLUE)

add_notes(slide, """这一页回答一个自然的问题：为什么 SCX 最早是在机器学习势函数这个领域做出来的，而不是在图像分类或者 NLP？

三个原因。

第一，ACE descriptor 本身就是一个天然的状态编码器。它把任意原子构型映射到一个 182 维的基函数空间，构型越相似，基函数向量越接近。你不用像做图像那样去训练一个 ResNet 来提特征——物理本身给了你编码器。

第二，DFT 提供了可靠但昂贵的参考标签。你可以把 DFT 看作一个"完美但很慢的专家"。它给你的标签一般是可靠的，所以你可以在上面画 residual map，判断你的势函数在哪个构型区域预测得好、哪个区域预测得差。这是有 ground truth 的。

第三，MLIP 天然就是多专家场景。同一个训练集可以训出多个势函数——不同架构、不同初始化、不同数据子集。它们在同一个构型区域的表现天然不同。SCX 要做的就是判断：在这个构型区域，哪个势函数最可靠？

这三个条件——天然编码器、可靠标签、多专家对比——让 MLIP 成了 SCX 最理想的第一个验证场景。""")

# ═══════════════════════════════════════════════
# SLIDE 7: SCX-MLIP Current Work
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_bg(slide, WHITE)
add_thin_line(slide, 0, 0, 13.333)

add_textbox(slide, 0.8, 0.5, 11.7, 0.7, "SCX-MLIP：第一个根据地的进展", font_size=28, bold=True, color=DARK)

# Left: ACE/Model B side
add_rect(slide, 0.5, 1.5, 6.0, 5.3, fill_color=LIGHT_GRAY)
add_textbox(slide, 0.8, 1.6, 5.4, 0.4, "势函数侧：Model B + Gauge Fixing", font_size=17, bold=True, color=DARK)
add_multiline(slide, 0.8, 2.1, 5.4, 4.5, [
    "• Single ACE (1196参数): 力RMSE 0.045 eV/Å",
    "",
    "• Model B shared+correction (1378参数):",
    "  E = Σ[c₀ᵀBᵢ + c_ZᵀBᵢ + b_Z]",
    "  shared c₀ + element correction c_Z",
    "",
    "• Gauge Fixing: post-hoc projection",
    "  训练后正交投影, violation 10⁻¹⁶",
    "  C33: +25.9% → +7.4% (vs DFT)",
    "  C44: -12.4% → -3.9% (vs DFT)",
    "",
    "• 关键认识: gauge fix 是势函数合并的",
    "  数学基础; 单独 AlN 不带来额外精度",
], font_size=12, color=GRAY)

# Right: SCX side
add_rect(slide, 6.8, 1.5, 6.0, 5.3, fill_color=LIGHT_BLUE)
add_textbox(slide, 7.1, 1.6, 5.4, 0.4, "SCX 侧：两层描述符 + 数据防中毒", font_size=17, bold=True, color=DARK)
add_multiline(slide, 7.1, 2.1, 5.4, 4.5, [
    "• 一层(12-dim MLIP encoder):",
    "  50%帧挤在一个状态, 噪声F1=0.253",
    "",
    "• 两层(=一层+ErrorDrivenEncoder):",
    "  L1粗聚类 → L2误差子空间细化",
    "  噪声F1 0.253→0.585 (+131%)",
    "  Top-3捕获94.6%噪声帧",
    "  Phonon 100%隔离到低误差状态",
    "",
    "• 数据防中毒 (Data Poisoning Defense):",
    "  无先验知识发现全部74个高噪声帧",
    "  训练fmax vs 测试误差 r=0.966",
    "  预估去噪后力RMSE降 29-48%",
], font_size=12, color=GRAY)

# Bottom highlight boxes
add_rect(slide, 0.5, 6.9, 4.0, 0.5, fill_color=BLUE)
add_textbox(slide, 0.7, 6.92, 3.6, 0.45, "534帧 → 可压缩~150帧 + 移除14帧极端噪声", font_size=13, bold=True, color=WHITE)

add_rect(slide, 4.8, 6.9, 3.5, 0.5, fill_color=RED_ACCENT)
add_textbox(slide, 5.0, 6.92, 3.1, 0.45, "SCX 100%命中 REPORT 最差12帧", font_size=13, bold=True, color=WHITE)

add_rect(slide, 8.6, 6.9, 4.2, 0.5, fill_color=ACCENT)
add_textbox(slide, 8.8, 6.92, 3.8, 0.45, "预估力RMSE: 0.045 → 0.023-0.032 eV/Å", font_size=13, bold=True, color=WHITE)

add_notes(slide, """这一页是第一个根据地的详细进展，分两个部分。

左边是势函数的工作。我们从 Single ACE 出发，做了 Model B 的 shared+correction 架构。这个架构把一个原子能分解为"所有原子共享的部分"加上"元素特异性修正"。但问题是这个分解不唯一——有个数学上的 gauge 自由度。我们后来发现用 post-hoc projection 可以在训练后无成本地消除这个自由度，violation 精确到 10 的负 16 次方。

右边是 SCX 的工作。我们用一层描述符分析 AlN 数据，发现 12 维手工特征太粗了——一半的帧挤在一个状态里，噪声检测的 F1 只有 0.253。然后我们做了两层方法：第一层粗聚类，第二层在误差相关的特征子空间里再细化。噪声 F1 直接翻了一倍多到了 0.585，三层状态就捕获了 94.6% 的噪声帧。

最让我兴奋的是"数据防中毒"这个发现——SCX 在不知道 fmax 阈值、不知道 batch 标签的情况下，独立发现了全部 74 个高噪声帧。而且训练时的 fmax 和测试时的预测误差的相关系数高达 0.966。这意味着去除 SCX 标记的噪声帧，测试集上的力 RMSE 预估能降 29% 到 48%。

底部三个数字是核心证据：534 帧中 SCX 建议压缩约 150 帧、移除 14 帧极端异常；噪声状态 100% 命中 REPORT 的最差 12 帧；预估力 RMSE 从 0.045 降到 0.023-0.032。""")

# ═══════════════════════════════════════════════
# SLIDE 8: From MLIP to SCX-Sim
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_bg(slide, WHITE)
add_thin_line(slide, 0, 0, 13.333)

add_textbox(slide, 0.8, 0.5, 11.7, 0.7, "从 MLIP 到科学仿真：问题结构是一样的", font_size=28, bold=True, color=DARK)

# Mapping table
add_multiline(slide, 0.8, 1.5, 11.7, 0.8, [
    "核心逻辑：状态 = 仿真参数空间中的区域 | 专家 = 不同保真度的求解器 | 数据价值 = 是否值得在这个参数区域做高保真计算",
], font_size=16, color=BLUE)

# Three column mapping
mappings = [
    ("MLIP (已验证)", [
        "状态 = 原子构型空间",
        "专家 = 不同势函数",
        "高保真 = DFT 计算",
        "低保真 = 经验势/ACE",
        "噪声 = fmax 异常的帧",
        "冗余 = 过多相似构型",
    ]),
    ("FEM/工程仿真", [
        "状态 = 参数空间(载荷/边界)",
        "专家 = 粗/细网格求解器",
        "高保真 = 精细网格 FEM",
        "低保真 = 粗网格/降阶模型",
        "噪声 = 数值发散/不收敛",
        "冗余 = 参数扫描过密区",
    ]),
    ("通用 SCX-Sim", [
        "状态 = 仿真场景参数组合",
        "专家 = 各保真度求解器链",
        "高保真 = 实验/高精度仿真",
        "低保真 = surrogate/经验公式",
        "噪声 = 收敛失败/数值异常",
        "冗余 = 参数空间过度采样",
    ]),
]

for i, (title, items) in enumerate(mappings):
    x = 0.5 + i * 4.2
    add_rect(slide, x, 2.6, 3.9, 4.0, fill_color=LIGHT_GRAY if i != 0 else LIGHT_BLUE)
    add_textbox(slide, x + 0.2, 2.7, 3.5, 0.4, title, font_size=16, bold=True, color=DARK)
    add_multiline(slide, x + 0.2, 3.2, 3.5, 3.2, items, font_size=13, color=GRAY if i != 0 else DARK)

# Bottom
add_textbox(slide, 0.8, 6.8, 11.7, 0.5,
    "结论：MLIP 中发现的规律（噪声集中在极端参数区、冗余集中在低误差区、专家在特定状态差异大）在 FEM/工程仿真中同样成立",
    font_size=15, bold=False, color=DARK)

add_notes(slide, """这一页是想说明：MLIP 不是 SCX 的终点，而是一个验证案例。

仔细看科学和工程仿真里的问题结构，和 MLIP 是完全一样的：

状态是什么？在 MLIP 里是原子构型空间，在 FEM 里是载荷、边界条件、材料参数组合的参数空间。

专家是什么？在 MLIP 里是不同的势函数，在仿真里是不同保真度的求解器——粗网格 FEM、细网格 FEM、降阶模型、实验数据。高保真计算很贵，低保真计算便宜但不太准。

噪声是什么？在 MLIP 里是 fmax 异常的帧，在仿真里是数值发散、收敛失败、网格畸变的点。

冗余是什么？在 MLIP 里是 phonon 小位移的 120 个相似帧，在仿真里是参数空间扫描过密的区域。

问题结构完全一样。SCX 在 MLIP 上发现的规律——噪声集中在极端参数区、冗余集中在低误差区、专家在特定状态差异大——在工程仿真里同样成立。这就是从 SCX-MLIP 到 SCX-Sim 的推广逻辑。""")

# ═══════════════════════════════════════════════
# SLIDE 9: SCX Workflow
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_bg(slide, WHITE)
add_thin_line(slide, 0, 0, 13.333)

add_textbox(slide, 0.8, 0.5, 11.7, 0.7, "SCX 通用工作流：五阶段闭环", font_size=28, bold=True, color=DARK)

# Flow diagram as labeled boxes with arrows
stages = [
    ("① 编码", "数据 → 特征向量\n(ACE/CNN/GNN/\n仿真参数向量)", BLUE),
    ("② 建图", "聚类 → 状态地图\n(KMeans/谱聚类/\n误差驱动聚类)", RGBColor(0x25, 0x69, 0xEB)),
    ("③ 估值", "四维评分\n质量·冗余·噪声\n·专家可靠性", ACCENT),
    ("④ 决策", "动作选择\n采集/压缩/重标\n/路由/丢弃", ORANGE),
    ("⑤ 反馈", "真实结果回传\n重校准阈值\n持续更新", RGBColor(0x7C, 0x3A, 0xED)),
]

for i, (title, desc, color) in enumerate(stages):
    x = 0.3 + i * 2.6
    # Box
    shape = add_rect(slide, x, 1.6, 2.3, 3.0, fill_color=WHITE, border_color=color)
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.6), Inches(2.3), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_textbox(slide, x + 0.1, 1.63, 2.1, 0.5, title, font_size=18, bold=True, color=WHITE)
    # Description
    add_multiline(slide, x + 0.15, 2.3, 2.0, 2.1, desc.split('\n'), font_size=12, color=DARK)

    # Arrow between boxes (except last)
    if i < 4:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.35), Inches(2.8), Inches(0.22), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color
        arrow.line.fill.background()

# Feedback loop arrow at bottom
add_textbox(slide, 0.8, 5.0, 11.7, 0.4,
    "反馈闭环：⑤ 的真实结果（训练loss变化、仿真验证误差）→ 回传校准 ③ 的阈值 → 持续改进决策质量",
    font_size=14, color=GRAY)

# Bottom: two-layer detail
add_rect(slide, 0.5, 5.6, 12.3, 1.7, fill_color=LIGHT_GRAY)
add_textbox(slide, 0.8, 5.7, 11.7, 0.4, "两层描述符（核心创新）", font_size=16, bold=True, color=DARK)
add_multiline(slide, 0.8, 6.1, 5.5, 1.0, [
    "Layer 1（人工域知识）：手工描述符 → 粗聚类",
    "  MLIP: coord_num, bond_len, volume...",
    "  仿真: Re, Ma, 载荷范围, 几何参数...",
], font_size=12, color=GRAY)
add_multiline(slide, 6.8, 6.1, 5.5, 1.0, [
    "Layer 2（算法发现）：误差驱动特征选择 → 细化",
    "  互信息筛选误差相关维度",
    "  仅在误差子空间中聚类",
    "  自动发现 noise/value/redundant 子状态",
], font_size=12, color=BLUE)

add_notes(slide, """这一页是 SCX 的通用流程图。

整个框架分成五个阶段，形成闭环。

第一阶段是编码——把不管你是什么类型的数据（原子构型、图像、仿真参数），都映射成特征向量。不同的领域用不同的编码器，但接口是统一的。

第二阶段是建图——在特征空间里聚类，形成状态地图。这里的关键创新是两层方法：第一层用人工域知识做粗聚类，第二层用误差驱动的特征选择做细化。这个两层设计是 SCX 区分噪声和困难样本的关键。

第三阶段是估值——对每个状态做四维评分：质量好不好、是不是冗余、有没有噪声风险、哪个专家最可靠。

第四阶段是决策——根据评分给出动作建议：继续采集、压缩、重新标注、路由给专家、或者直接丢弃。

第五阶段是反馈——真实的训练结果或仿真验证误差回传过来，校准第三阶段的阈值。这是一个持续改进的闭环。

底部的两层描述符是 SCX 最大的方法论创新。这个想法的本质是：状态应该由"模型在哪里失败"来定义，而不是由人的直觉来定义。""")

# ═══════════════════════════════════════════════
# SLIDE 10: Application Directions
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_bg(slide, WHITE)
add_thin_line(slide, 0, 0, 13.333)

add_textbox(slide, 0.8, 0.5, 11.7, 0.7, "SCX 的潜在应用方向：同一个问题结构", font_size=28, bold=True, color=DARK)

add_textbox(slide, 0.8, 1.2, 11.7, 0.4,
    '核心问题（普适的）："在当前状态 s 下，这个数据/模型是否可信？"',
    font_size=16, color=BLUE)

apps = [
    ("材料/DFT/MLIP ✅", "已验证", "势函数选择、数据质量筛查、训练集压缩", BLUE),
    ("FEM/工程仿真", "结构相同", "粗/细网格切换策略、仿真点冗余检测、数值异常筛查", ACCENT),
    ("制造在线监测", "结构相同", "传感器数据质控、异常检测(非统计型)、多模型路由", RGBColor(0x25, 0x69, 0xEB)),
    ("医学图像质控", "结构相同", "标注质量评估、困难样本 vs 标注错误区分、冗余切片筛选", ORANGE),
    ("机器人轨迹数据", "结构相同", "仿真轨迹质量评估、sim-to-real 场景选择、多控制器路由", RGBColor(0x7C, 0x3A, 0xED)),
    ("智能驾驶场景", "结构相同", "Corner case 价值评估、仿真场景冗余检测、多传感器融合路由", RED_ACCENT),
]

for i, (name, status, detail, color) in enumerate(apps):
    y = 1.8 + i * 0.85
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.07), Inches(0.65))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_textbox(slide, 0.8, y + 0.02, 2.8, 0.6, name, font_size=16, bold=True, color=DARK)
    add_textbox(slide, 3.8, y + 0.02, 1.5, 0.6, status, font_size=13, color=color)
    add_textbox(slide, 5.5, y + 0.02, 7.3, 0.6, detail, font_size=14, color=GRAY)

# Disclaimer
add_textbox(slide, 0.8, 6.9, 11.7, 0.5,
    "坦诚声明：除 MLIP 外，其他方向尚未验证。列出是为了说明'问题结构相同'，不是宣称已经解决。",
    font_size=13, bold=False, color=RED_ACCENT)

add_notes(slide, """这一页列出了 SCX 可以应用的几个方向，但我必须坦诚地说——除了 MLIP 之外，其他方向都还没有实际验证。

列出它们的目的，是想说明一个问题：这些看似不相关的领域，底层的"数据价值判断"逻辑是相同的。都是问同一个问题："在当前这个状态下，这个数据/这个模型可不可信？"

材料计算是已经验证过的。FEM 仿真里，粗网格和细网格的切换策略、哪些仿真点已经算够了、哪些参数组合的收敛失败是真的数值问题——这些问题结构和 MLIP 里是一模一样的。

制造监测里的传感器数据质控——传统方法用统计阈值，但 SCX 可以考虑到"在当前工况下这个传感器的读数是否正常"，这是一个状态条件的判断。

医学图像的标注质量——放射科医生标注一张 CT 的时候，某些类型的病变天然就容易有分歧。SCX 可以区分"这个样本医生们意见不一致是因为它难"还是"这个样本标注确实错了"。

再次强调：这些只是潜在方向。我们现在扎实做好的只有 MLIP 这一个。""")

# ═══════════════════════════════════════════════
# SLIDE 11: Paper Value
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_bg(slide, WHITE)
add_thin_line(slide, 0, 0, 13.333)

add_textbox(slide, 0.8, 0.5, 11.7, 0.7, "为什么 SCX 值得写成论文", font_size=28, bold=True, color=DARK)

# Four pillars
pillars = [
    ("理论新意", "数据价值是状态条件的", [
        "不是全局标量，是条件概率",
        "Proposition 1 (全局排名不充分性) 已证明",
        "6 个命题覆盖压缩/路由/两层发现",
        "重新定义了"数据价值"的数学对象",
    ], BLUE),
    ("应用验证", "MLIP 场景已有量化证据", [
        "AlN v3: 噪声F1=0.585, Top-3 捕获94.6%",
        "数据防中毒: 100%命中, r=0.966",
        "MedMNIST compress: 30%压缩精度+6%",
        "MedMNIST routing: SCX > ensemble",
    ], ACCENT),
    ("通用性", "不限于 ML/DL 范式", [
        "不依赖特定模型架构",
        "不依赖梯度/损失函数形式",
        "抽象基类 + YAML 域配置",
        "3 步接入新领域",
    ], ORANGE),
    ("实用性", "直接回答"值不值得花钱"", [
        "数据采集: 哪些区域值得继续算",
        "数据清洗: 自动发现噪声/异常",
        "专家部署: 路由到最可靠的模型",
        "成本优化: 压缩冗余, 聚焦高价值",
    ], RGBColor(0x7C, 0x3A, 0xED)),
]

for i, (title, subtitle, points, color) in enumerate(pillars):
    x = 0.3 + i * 3.25
    add_rect(slide, x, 1.5, 3.0, 5.2, fill_color=LIGHT_GRAY)
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.5), Inches(3.0), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_textbox(slide, x + 0.15, 1.52, 2.7, 0.4, title, font_size=18, bold=True, color=WHITE)
    add_textbox(slide, x + 0.15, 1.92, 2.7, 0.3, subtitle, font_size=11, color=WHITE)
    # Points
    add_multiline(slide, x + 0.15, 2.5, 2.7, 4.0, points, font_size=12, color=DARK if i == 0 else GRAY)

add_notes(slide, """这一页论证 SCX 的论文价值，从四个角度。

理论新意：SCX 重新定义了"数据价值"这个概念——它不是一个全局标量，而是一个状态条件概率。我们已经证明了 6 个数学命题，包括全局排名不充分性、压缩保真度下界、以及两层状态发现的形式化。这个理论框架是新的。

应用验证：在 MLIP 场景，我们有量化证据。AlN v3 上两层方法的噪声 F1 达到 0.585，数据防中毒 100% 命中。MedMNIST 上压缩实验证明 SCX 优于随机采样。路由实验证明状态条件权重优于均匀集成。这些不是概念演示，是有数字的。

通用性：SCX 不绑定任何特定的模型架构。不管是 ACE 线性模型还是 ResNet 还是 Transformer，只要你能定义状态编码器和误差信号，SCX 就能工作。我们已经用抽象基类 + YAML 配置实现了新领域 3 步接入。

实用性：SCX 直接回答一个很实际的问题——"这个数据值不值得花钱？"。不是抽象的"数据质量"，而是可操作的建议：继续采集、压缩、重标、还是路由。""")

# ═══════════════════════════════════════════════
# SLIDE 12: Paper Roadmap
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_bg(slide, WHITE)
add_thin_line(slide, 0, 0, 13.333)

add_textbox(slide, 0.8, 0.5, 11.7, 0.7, "计划中的两篇论文", font_size=28, bold=True, color=DARK)

# Paper 1
add_rect(slide, 0.5, 1.5, 5.8, 4.5, fill_color=LIGHT_BLUE, border_color=BLUE)
add_textbox(slide, 0.8, 1.6, 5.2, 0.5, "Paper 1: SCX-MLIP", font_size=22, bold=True, color=BLUE)
add_textbox(slide, 0.8, 2.15, 5.2, 0.4, "专家势函数合并与状态条件数据价值", font_size=16, color=DARK)
add_multiline(slide, 0.8, 2.7, 5.2, 3.0, [
    "核心贡献：",
    "• Model B shared+correction 架构",
    "• Post-hoc gauge fixing 方法",
    "• SCX 两层噪声/冗余检测 (AlN v3)",
    "• 数据防中毒 (data poisoning defense)",
    "",
    "实验基础：",
    "• AlN v3: 534帧 DFT, 689 清洁标签",
    "• Single ACE vs Model B 完整对比",
    "• SCX 噪声检测 100% 命中验证",
    "",
    "目标期刊：Nature Comp Sci / npj CM",
], font_size=13, color=GRAY)

# Paper 2
add_rect(slide, 6.8, 1.5, 5.8, 4.5, fill_color=LIGHT_GRAY, border_color=ACCENT)
add_textbox(slide, 7.1, 1.6, 5.2, 0.5, "Paper 2: SCX-Theory", font_size=22, bold=True, color=ACCENT)
add_textbox(slide, 7.1, 2.15, 5.2, 0.4, "The value of data is state-conditioned", font_size=16, color=DARK)
add_multiline(slide, 7.1, 2.7, 5.2, 3.0, [
    "核心贡献：",
    "• 通用数学框架 (6 个命题)",
    "• 两层描述符方法论",
    "• 数据四分类理论",
    "• 状态条件专家路由",
    "",
    "实验基础：",
    "• MLIP (已验证) + MedMNIST (已验证)",
    "• 合成数据 benchmark",
    "• 计划: FEM 仿真案例",
    "",
    "目标期刊：TMLR / AISTATS",
], font_size=13, color=GRAY)

# Timeline
add_textbox(slide, 0.8, 6.3, 11.7, 0.4, "写作节奏", font_size=18, bold=True, color=DARK)
add_multiline(slide, 0.8, 6.7, 11.7, 0.6, [
    "Paper 1 可先行（已有完整实验数据）→ Paper 2 需要额外完成 FEM 仿真案例 + 合成 benchmark 后投稿",
], font_size=14, color=GRAY)

add_notes(slide, """计划写两篇论文。

第一篇 SCX-MLIP，聚焦在势函数合并和数据质量。这已经有完整的实验数据了——689 帧 DFT 标签、Single ACE vs Model B 的完整对比、SCX 两层方法在 AlN v3 上的验证，还有数据防中毒的 100% 命中。这篇的叙事线很清楚：从 gauge fixing 到状态条件数据管理，有方法论也有实验。

第二篇 SCX-Theory，把框架推广到通用科学仿真。这需要更多的理论工作和额外的实验验证——FEM 仿真案例、合成数据 benchmark。但这篇的定位更高：不是在说"MLIP 数据怎么管"，而是在说"数据价值的定义应该重新思考"。

时间线上，Paper 1 现在就可以开始写。Paper 2 等 FEM 案例和 benchmark 做完再投稿。""")

# ═══════════════════════════════════════════════
# SLIDE 13: Closing
# ═══════════════════════════════════════════════
slide = add_blank_slide()
add_bg(slide, WHITE)

# Large closing statement
add_textbox(slide, 1.0, 2.0, 11.3, 1.5,
    "SCX 不替代模型，也不替代专家。",
    font_size=36, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.0, 3.2, 11.3, 0.8,
    "SCX 判断：在当前状态下，",
    font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.0, 3.9, 11.3, 1.0,
    "哪些数据值得信、哪些专家值得信、哪些动作值得花钱。",
    font_size=28, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)

# Divider
add_thin_line(slide, 4.5, 5.2, 4.3)

add_textbox(slide, 1.0, 5.6, 11.3, 0.6,
    "The value of data is state-conditioned.",
    font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.0, 6.3, 11.3, 0.5,
    "谢谢各位老师 | 欢迎提问", font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

add_notes(slide, """最后一页，一句话总结。

SCX 不替代任何模型，也不替代任何专家。它做的事情很朴素——判断在当前这个状态下，哪些数据值得你信、哪些专家值得你信、哪些动作值得你花钱。

这不是一个"更好的模型"，而是一个"更好的决策框架"。它的核心命题只有一句话：The value of data is state-conditioned。数据价值是状态条件的。

谢谢各位老师，欢迎提问。""")

# ── Save ──
output_path = "G:/Xiaogan_Supercomputing_data/SCX/outputs/SCX_学术汇报_2026-06-26.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
