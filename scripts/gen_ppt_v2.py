#!/usr/bin/env python
"""Generate SCX academic PPTX -- clean version."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x1A, 0x56, 0xDB)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x0D, 0x94, 0x8B)
RED_ACCENT = RGBColor(0xDC, 0x26, 0x26)
ORANGE = RGBColor(0xE6, 0x8A, 0x00)
LIGHT_BLUE = RGBColor(0xDB, 0xE9, 0xFA)
HEADER_BG = RGBColor(0x25, 0x69, 0xEB)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PINK_BG = RGBColor(0xFD, 0xE8, 0xE8)
WARM_BG = RGBColor(0xFE, 0xF3, 0xE6)

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_bg(slide, color=LIGHT_GRAY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def hline(slide, l, t, w):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Pt(3))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()

def tb(slide, l, t, w, h, text, fs=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(fs)
    p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    return tf

def ml(slide, l, t, w, h, lines, fs=16, color=DARK):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(fs); p.font.color.rgb = color; p.space_after = Pt(4)
    return tf

def card(slide, l, t, w, h, fill=WHITE, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def bar(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ====== SLIDE 1: Cover ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 1.5, 1.8, 10.3, 1.2, 'SCX: State-Conditioned Data Value Framework', fs=40, bold=True, color=DARK)
tb(s, 1.5, 3.0, 10.3, 0.8, 'The value of data is state-conditioned', fs=24, color=BLUE)
tb(s, 1.5, 3.8, 10.3, 0.5, 'From ML Potentials to Scientific & Engineering Simulation', fs=18, color=GRAY)
tb(s, 1.5, 5.5, 5, 1.0, 'Speaker: Shaojie Xu\nHuazhong Univ. of Sci. & Tech.\nJune 2026', fs=16, color=GRAY)
notes(s, 'Good morning. Today I present SCX -- a framework for state-conditioned data valuation. The core idea: data value is not a fixed property, but depends on state context. I will cover four things: the problem, the core idea, what we have built, and next steps.')

# ====== SLIDE 2: Problem ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'Problem: Data is expensive, but we often do not know which data is worth using', fs=26, bold=True, color=DARK)
problems = [
    ('Data Acquisition\nis Costly', ['DFT: tens of core-hours per structure', 'FEM: hours per simulation run', 'Medical imaging: expert annotation needed']),
    ('Data Value\nis Unclear', ['Which samples are worth computing?', 'Which are redundant / can be skipped?', 'Which may be wrong?']),
    ('Multi-Expert Comparison\nis Hard', ['No single best model exists globally', 'Expert reliability varies by region', 'How to route samples to the right expert?']),
    ('Traditional Methods\nFall Short', ['LOO/Shapley: one fixed score per sample', 'No state context considered', 'Cannot separate noise from hard samples']),
]
for i, (title, items) in enumerate(problems):
    x = 0.5 + i * 3.1
    card(s, x, 1.5, 2.9, 3.8, fill=LIGHT_GRAY)
    tb(s, x + 0.2, 1.6, 2.5, 0.7, title, fs=15, bold=True, color=DARK)
    ml(s, x + 0.2, 2.4, 2.5, 2.8, items, fs=13, color=GRAY)
tb(s, 0.8, 5.7, 11.7, 1.0, 'Core tension: We need a framework that judges -- at the state level -- whether this data is trustworthy, this expert is reliable, and this action is worth the cost.', fs=17, color=DARK)
notes(s, 'Data is expensive in scientific computing. DFT calculations cost dozens of core-hours. But we often do not know which data to invest in. Traditional methods give each sample a fixed score, ignoring context. This is the gap SCX fills.')

# ====== SLIDE 3: Core Idea ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'Core Idea: Data Value is State-Conditioned', fs=28, bold=True, color=DARK)
card(s, 0.5, 1.5, 5.8, 2.0, fill=LIGHT_GRAY)
tb(s, 0.8, 1.6, 5.2, 0.4, 'Traditional: Global Scoring', fs=20, bold=True, color=RED_ACCENT)
ml(s, 0.8, 2.1, 5.2, 1.2, ['LOO / Data Shapley / Influence Function', 'Each sample -> one fixed score', 'Ignores: who uses it, what exists, what is missing', 'Result: noise vs hard sample indistinguishable'], fs=14, color=GRAY)
card(s, 6.8, 1.5, 5.8, 2.0, fill=LIGHT_BLUE)
tb(s, 7.1, 1.6, 5.2, 0.4, 'SCX: State-Conditioned Scoring', fs=20, bold=True, color=BLUE)
ml(s, 7.1, 2.1, 5.2, 1.2, ['V(x) = V(x | state, coverage, expert, cost, feedback)', 'Same sample -> different value under different conditions', 'Data value determined by five factors jointly', 'Can separate: redundant / noisy / valuable / expert-dep.'], fs=14, color=DARK)
card(s, 0.5, 3.9, 12.3, 1.2, fill=WHITE, border=BLUE)
tb(s, 0.8, 4.0, 11.7, 0.4, 'Core Formula', fs=16, bold=True, color=BLUE)
tb(s, 0.8, 4.4, 11.7, 0.5, 'V(s) = r_bar(s) * rho(s) * L(s) * [1 - D(s)] * max_m SCX_m(s)', fs=22, bold=True, color=DARK)
tb(s, 0.8, 4.8, 11.7, 0.3, 'r_bar=current error | rho=state frequency | L=learnability (non-noise) | D=coverage | SCX=best expert reliability in state s', fs=11, color=GRAY)
tb(s, 0.8, 5.5, 11.7, 1.0, 'Key insight: Data value is not an intrinsic property. The same data can be noise to Model A but high-value to Model B.', fs=17, color=DARK)
notes(s, 'Core idea. Traditional methods assign fixed scores. SCX argues data value depends on state. The formula V(s) has five multiplicative factors. If any one is zero, value is zero. For example, if current error is already low, more data is redundant regardless of sample quality.')

# ====== SLIDE 4: Four-quadrant with AlN data ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'Plain Language: Same Sample, Different Value in Different States', fs=26, bold=True, color=DARK)
quads = [
    ('Redundant', 'Low error + High density', 'Compress / Skip', 'Phonon 120 frames: fmax~0.66\nEstimated 50-80% compressible', LIGHT_GRAY),
    ('Valuable', 'High error + High density\n+ High consistency', 'Acquire More', 'EOS extreme-strain region\nElastic deformation directions', LIGHT_BLUE),
    ('Noisy / Bad Labels', 'High error + Low density\n+ Low consistency', 'Recompute/Downweight', 'Thermal 1800K: 49.1% fmax>5\n14 frames fmax>10 (extreme anomaly)', PINK_BG),
    ('Expert-Dependent', 'One expert significantly\nbetter in this state', 'Route to Best Expert', 'MLMD stress10 region\nRoute to high-T specialist', WARM_BG),
]
for i, (title, cond, action, example, bg) in enumerate(quads):
    x = 0.3 + (i % 2) * 6.5
    y = 1.4 + (i // 2) * 2.8
    card(s, x, y, 6.2, 2.5, fill=bg)
    tb(s, x + 0.2, y + 0.1, 5.8, 0.4, title, fs=18, bold=True, color=DARK)
    tb(s, x + 0.2, y + 0.5, 5.8, 0.4, cond, fs=13, color=BLUE)
    tb(s, x + 0.2, y + 0.95, 5.8, 0.3, '-> ' + action, fs=14, bold=True, color=DARK)
    ml(s, x + 0.2, y + 1.3, 5.8, 1.1, example.split('\n'), fs=12, color=GRAY)
notes(s, 'Four-quadrant explanation with real AlN data. Redundant: phonon 120 frames with low fmax, model already learned -- compress 50-80%. Valuable: EOS extreme strain with higher error but good consistency -- real physics not yet learned. Noisy: thermal 1800K with 49% frames above fmax=5 -- these are bad VASP configurations. Expert-dependent: MLMD region needs routing to the right specialist.')

# ====== SLIDE 5: Four-classification ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'SCX Four-Class Data Taxonomy', fs=28, bold=True, color=DARK)
classes = [
    ('Valuable', 'High err + High density + High consistency + Low coverage', 'Acquire More', 'EOS extreme strain, elastic deformation regions', BLUE),
    ('Redundant', 'Low err + High density + High coverage', 'Compress / Skip', 'Phonon 120 frames: compress 50-80% without accuracy loss', ACCENT),
    ('Noisy-Risk', 'High err + Low density + Low consistency', 'Recompute / Downweight / Discard', 'Thermal 1800K: 49% anomaly; MLMD stress10: 35% anomaly', RED_ACCENT),
    ('Expert-Dep.', 'One expert significantly better than others here', 'Route to Specialist', 'MLMD high-T needs specialist; surface/defect needs transfer', ORANGE),
]
for i, (name, cond, action, example, clr) in enumerate(classes):
    y = 1.4 + i * 1.45
    bar(s, 0.5, y, 0.08, 1.2, clr)
    tb(s, 0.8, y + 0.05, 2.0, 1.1, name, fs=18, bold=True, color=DARK)
    tb(s, 3.0, y + 0.05, 3.0, 1.1, cond, fs=14, color=GRAY)
    tb(s, 6.2, y + 0.05, 2.5, 1.1, '-> ' + action, fs=16, bold=True, color=clr)
    tb(s, 9.0, y + 0.05, 4.0, 1.1, example, fs=13, color=GRAY)
notes(s, 'SCX classifies data into four types with clear conditions and actions. Valuable: model has not learned this yet, data quality is reliable -- invest more. Redundant: model already performs well -- save money. Noisy: data itself is problematic -- do not use. Expert-dependent: route to the right specialist, do not average.')

# ====== SLIDE 6: Why MLIP first ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'Why MLIP is the First Proving Ground for SCX', fs=28, bold=True, color=DARK)
reasons = [
    ('1. Natural State Encoder', 'ACE descriptor maps any atomic configuration to a 182-dim basis space. Similar configurations -> similar basis vectors. No need to train an extra embedding network. Physics provides the encoder.'),
    ('2. Reliable Reference Labels', 'DFT provides trustworthy (though expensive) energy/force/stress labels. With ground truth, we can draw residual maps and judge expert quality. This provides an anchor absent in purely unsupervised settings.'),
    ('3. Natural Multi-Expert Scenario', 'Same training set -> multiple potentials (different architectures, initializations, data subsets). Their performance naturally differs across configuration regions -- exactly what SCX needs to compare.'),
]
for i, (title, body) in enumerate(reasons):
    y = 1.5 + i * 1.6
    card(s, 0.5, y, 12.3, 1.4, fill=LIGHT_GRAY if i % 2 == 0 else LIGHT_BLUE)
    tb(s, 0.8, y + 0.1, 11.7, 0.4, title, fs=20, bold=True, color=DARK)
    tb(s, 0.8, y + 0.55, 11.7, 0.7, body, fs=14, color=GRAY)
tb(s, 0.8, 6.5, 11.7, 0.5, 'Summary: MLIP = Natural Encoder (ACE) + Reliable Labels (DFT) + Multi-Expert Comparison -> Ideal Starting Point for SCX', fs=17, bold=True, color=BLUE)
notes(s, 'Why MLIP first? Three reasons. ACE descriptor is a natural state encoder -- physics gives us the embedding. DFT provides reliable labels -- we have ground truth. MLIP is naturally multi-expert -- different potentials on same data. These three conditions make MLIP the ideal first validation domain.')

# ====== SLIDE 7: SCX-MLIP Progress ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'SCX-MLIP: Progress at the First Proving Ground', fs=28, bold=True, color=DARK)
card(s, 0.5, 1.5, 6.0, 5.3, fill=LIGHT_GRAY)
tb(s, 0.8, 1.6, 5.4, 0.4, 'Potential Side: Model B + Gauge Fixing', fs=17, bold=True, color=DARK)
ml(s, 0.8, 2.1, 5.4, 4.5, [
    'Single ACE (1196 params): Force RMSE 0.045 eV/A',
    '',
    'Model B shared+correction (1378 params):',
    '  E = sum[c0*Bi + c_Z*Bi + b_Z]',
    '  shared c0 + element correction c_Z',
    '',
    'Gauge Fixing: post-hoc projection',
    '  Train first, then project: c_Z -= g, c_0 += g',
    '  Violation = 4.6e-16 (machine zero)',
    '  Zero prediction change (math identity)',
    '',
    'Key: gauge fixing is for potential MERGING',
    '(AlN+GaN). For pure AlN alone, it adds no accuracy.',
], fs=12, color=GRAY)
card(s, 6.8, 1.5, 6.0, 5.3, fill=LIGHT_BLUE)
tb(s, 7.1, 1.6, 5.4, 0.4, 'SCX Side: Two-Layer Descriptor + Defense', fs=17, bold=True, color=DARK)
ml(s, 7.1, 2.1, 5.4, 4.5, [
    'One-layer (12-dim MLIP encoder):',
    '  50% frames in one blob, Noise F1 = 0.253',
    '',
    'Two-layer (= L1 + ErrorDrivenEncoder):',
    '  L1 coarse -> L2 error-subspace refinement',
    '  Noise F1: 0.253 -> 0.585 (+131%)',
    '  Top-3 states capture 94.6% of noise frames',
    '  Phonon 120 frames 100% isolated (low-error)',
    '',
    'Data Poisoning Defense:',
    '  Found ALL 74 high-noise frames (no prior knowledge)',
    '  Train fmax vs test error: r = 0.966',
    '  Est. force RMSE improvement: 29-48%',
], fs=12, color=GRAY)
bar(s, 0.5, 6.9, 4.0, 0.5, BLUE); tb(s, 0.7, 6.92, 3.6, 0.4, '534 frames -> compress ~150 + remove 14 extreme', fs=13, bold=True, color=WHITE)
bar(s, 4.8, 6.9, 3.5, 0.5, RED_ACCENT); tb(s, 5.0, 6.92, 3.1, 0.4, 'SCX 100% hits REPORT worst-12 frames', fs=13, bold=True, color=WHITE)
bar(s, 8.6, 6.9, 4.2, 0.5, ACCENT); tb(s, 8.8, 6.92, 3.8, 0.4, 'Est. Force RMSE: 0.045 -> 0.023-0.032 eV/A', fs=13, bold=True, color=WHITE)
notes(s, 'Progress on two parallel tracks. Left: Model B architecture and gauge fixing via post-hoc projection. Right: SCX two-layer descriptor improving noise F1 2.3x, and data poisoning defense -- SCX found all high-noise frames without knowing fmax thresholds. Pearson r = 0.966 between training fmax and test error.')

# ====== SLIDE 8: MLIP -> SCX-Sim ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'From MLIP to Simulation: The Problem Structure is Identical', fs=28, bold=True, color=DARK)
tb(s, 0.8, 1.3, 11.7, 0.4, 'Core mapping: State = region in parameter space | Expert = solvers at different fidelities | Data value = is high-fidelity worth it here?', fs=15, color=BLUE)
mappings = [
    ('MLIP [Verified]', ['State = atomic config space', 'Expert = different potentials', 'High-fidelity = DFT calculation', 'Low-fidelity = empirical / ACE', 'Noise = frames with anomalous fmax', 'Redundancy = too many similar configs'], True),
    ('FEM / Engineering Sim', ['State = param space (load/BC)', 'Expert = coarse/fine mesh solvers', 'High-fidelity = fine-mesh FEM', 'Low-fidelity = coarse / ROM', 'Noise = divergence / non-convergence', 'Redundancy = over-sampled param regions'], False),
    ('Generic SCX-Sim', ['State = simulation scenario params', 'Expert = solver chain at each fidelity', 'High-fidelity = experiment / high-acc', 'Low-fidelity = surrogate / empirical', 'Noise = convergence failure / anomaly', 'Redundancy = over-sampled param space'], False),
]
for i, (title, items, is_verified) in enumerate(mappings):
    x = 0.5 + i * 4.2
    card(s, x, 2.0, 3.9, 4.0, fill=LIGHT_BLUE if is_verified else LIGHT_GRAY)
    tb(s, x + 0.2, 2.1, 3.5, 0.4, title, fs=15, bold=True, color=DARK)
    ml(s, x + 0.2, 2.6, 3.5, 3.2, items, fs=13, color=DARK if is_verified else GRAY)
tb(s, 0.8, 6.4, 11.7, 0.8, 'Patterns discovered in MLIP (noise in extreme regions, redundancy in low-error regions, expert divergence in specific states) also hold in FEM and engineering simulation.', fs=15, color=DARK)
notes(s, 'The problem structure is identical across MLIP and engineering simulation. State spaces, expert types, noise patterns, redundancy patterns -- all have direct analogs. The mapping is clean and suggests SCX generalizes naturally.')

# ====== SLIDE 9: Workflow ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'SCX General Workflow: Five-Stage Closed Loop', fs=28, bold=True, color=DARK)
stages = [
    ('1. Encode', 'Data -> Feature Vector\n(ACE/CNN/GNN/\nsim param vector)', BLUE),
    ('2. Map', 'Cluster -> State Map\n(KMeans/Spectral/\nerror-driven clustering)', HEADER_BG),
    ('3. Evaluate', '4-Dimension Scoring\nQuality/Redundancy/\nNoise/Expert Reliability', ACCENT),
    ('4. Decide', 'Action Selection\nAcquire/Compress/\nRelabel/Route/Discard', ORANGE),
    ('5. Feedback', 'Real Results Return\nRecalibrate Thresholds\nContinuous Update', PURPLE),
]
for i, (title, desc, clr) in enumerate(stages):
    x = 0.3 + i * 2.6
    card(s, x, 1.6, 2.3, 3.0, fill=WHITE, border=clr)
    bar(s, x, 1.6, 2.3, 0.55, clr)
    tb(s, x + 0.1, 1.63, 2.1, 0.5, title, fs=18, bold=True, color=WHITE)
    ml(s, x + 0.15, 2.3, 2.0, 2.1, desc.split('\n'), fs=12, color=DARK)
    if i < 4:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.35), Inches(2.8), Inches(0.22), Inches(0.4))
        arr.fill.solid(); arr.fill.fore_color.rgb = clr; arr.line.fill.background()
tb(s, 0.8, 5.0, 11.7, 0.4, 'Feedback loop: Stage 5 results -> recalibrate Stage 3 thresholds -> continuous improvement', fs=14, color=GRAY)
card(s, 0.5, 5.6, 12.3, 1.7, fill=LIGHT_GRAY)
tb(s, 0.8, 5.7, 11.7, 0.4, 'Two-Layer Descriptor (Key Innovation)', fs=16, bold=True, color=DARK)
ml(s, 0.8, 6.1, 5.5, 1.0, ['Layer 1 (Human Domain Knowledge): Manual descriptors -> coarse clustering', '  MLIP: coord_num, bond_len, volume...', '  Simulation: Re, Ma, load range, geometry params...'], fs=12, color=GRAY)
ml(s, 6.8, 6.1, 5.5, 1.0, ['Layer 2 (Algorithm-Discovered): Error-driven feature selection -> refinement', '  Mutual info screening of error-correlated dimensions', '  Cluster ONLY in error-correlated subspace', '  Auto-discover noise/value/redundant sub-states'], fs=12, color=BLUE)
notes(s, 'SCX workflow: five stages forming a closed loop. Key innovation: two-layer descriptor. Layer 1 uses human domain knowledge. Layer 2 uses error-driven feature selection. States should be defined by where models FAIL, not by human intuition.')

# ====== SLIDE 10: Applications ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'Potential Applications: Same Problem Structure Across Domains', fs=28, bold=True, color=DARK)
tb(s, 0.8, 1.2, 11.7, 0.4, 'Universal question: Given the current state s, is this data / model trustworthy?', fs=16, color=BLUE)
apps = [
    ('Materials/DFT/MLIP [Verified]', 'Potential selection, data quality screening, training set compression', BLUE),
    ('FEM / Engineering Simulation', 'Coarse/fine mesh switching, simulation point redundancy, numerical anomaly detection', ACCENT),
    ('Manufacturing Online Monitoring', 'Sensor data QC, anomaly detection (non-statistical), multi-model routing', HEADER_BG),
    ('Medical Image Quality Control', 'Annotation quality, hard-case vs label-error separation, redundant slice filtering', ORANGE),
    ('Robot Trajectory Data', 'Sim trajectory quality, sim-to-real scenario selection, multi-controller routing', PURPLE),
    ('Autonomous Driving Scenarios', 'Corner case valuation, simulation scenario redundancy, multi-sensor fusion routing', RED_ACCENT),
]
for i, (name, detail, clr) in enumerate(apps):
    y = 1.8 + i * 0.82
    bar(s, 0.5, y, 0.07, 0.62, clr)
    tb(s, 0.8, y + 0.02, 3.5, 0.55, name, fs=15, bold=True, color=DARK)
    tb(s, 4.5, y + 0.02, 8.3, 0.55, detail, fs=14, color=GRAY)
tb(s, 0.8, 6.9, 11.7, 0.4, 'Honest disclaimer: Except MLIP, all other directions are unverified. Listed to illustrate shared problem structure, not to claim solutions.', fs=13, color=RED_ACCENT)
notes(s, 'Potential application directions. Only MLIP is verified. Others share the same problem structure -- state-conditioned data/model trustworthiness -- but have not been tested. The core question is universal: in the current state, is this data trustworthy? Is this model reliable?')

# ====== SLIDE 11: Paper Value ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'Why SCX Deserves Publication', fs=28, bold=True, color=DARK)
pillars = [
    ('Theory', 'Data value is\nstate-conditioned', ['Not a global scalar, but a conditional probability', 'Prop 1 (global ranking insufficiency) proven', '6 propositions: compress/routing/two-layer', 'Redefines data value mathematically'], BLUE),
    ('Validation', 'MLIP has\nquantitative evidence', ['AlN v3: Noise F1=0.585, Top-3 captures 94.6%', 'Data poisoning defense: 100% hit, r=0.966', 'MedMNIST compress: +6% at 30% reduction', 'MedMNIST routing: SCX > ensemble'], ACCENT),
    ('Generality', 'Not tied to\nany ML paradigm', ['No dependency on model architecture', 'No dependency on gradient/loss form', 'Abstract base classes + YAML config', 'New domain integration in 3 steps'], ORANGE),
    ('Practicality', 'Directly answers\nis it worth the cost', ['Data acquisition: which regions to invest', 'Data cleaning: auto-discover noise', 'Expert deployment: route to best model', 'Cost optimization: compress, focus'], PURPLE),
]
for i, (title, subtitle, points, clr) in enumerate(pillars):
    x = 0.3 + i * 3.25
    card(s, x, 1.5, 3.0, 5.2, fill=LIGHT_GRAY)
    bar(s, x, 1.5, 3.0, 0.8, clr)
    tb(s, x + 0.15, 1.52, 2.7, 0.4, title, fs=18, bold=True, color=WHITE)
    tb(s, x + 0.15, 1.92, 2.7, 0.3, subtitle, fs=11, color=WHITE)
    ml(s, x + 0.15, 2.5, 2.7, 4.0, points, fs=12, color=DARK if i == 0 else GRAY)
notes(s, 'Four pillars: Theory (redefines data value as state-conditioned, 6 propositions), Validation (quantitative evidence from AlN and MedMNIST), Generality (not tied to any ML paradigm), Practicality (directly answers cost questions).')

# ====== SLIDE 12: Paper Roadmap ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'Planned Two-Paper Sequence', fs=28, bold=True, color=DARK)
card(s, 0.5, 1.5, 5.8, 4.5, fill=LIGHT_BLUE, border=BLUE)
tb(s, 0.8, 1.6, 5.2, 0.5, 'Paper 1: SCX-MLIP', fs=22, bold=True, color=BLUE)
tb(s, 0.8, 2.1, 5.2, 0.4, 'Expert Potential Merging & State-Conditioned Data Value', fs=15, color=DARK)
ml(s, 0.8, 2.7, 5.2, 3.0, [
    'Contributions:',
    '- Model B shared+correction architecture',
    '- Post-hoc gauge fixing method',
    '- SCX two-layer noise/redundancy detection',
    '- Data poisoning defense',
    '',
    'Experiments:',
    '- AlN v3: 534 DFT frames, 689 clean labels',
    '- Single ACE vs Model B full comparison',
    '- SCX 100% noise hit rate verified',
    '',
    'Target: Nature Comp Sci / npj Comp Mater',
], fs=13, color=GRAY)
card(s, 6.8, 1.5, 5.8, 4.5, fill=LIGHT_GRAY, border=ACCENT)
tb(s, 7.1, 1.6, 5.2, 0.5, 'Paper 2: SCX-Theory', fs=22, bold=True, color=ACCENT)
tb(s, 7.1, 2.1, 5.2, 0.4, 'The Value of Data is State-Conditioned', fs=15, color=DARK)
ml(s, 7.1, 2.7, 5.2, 3.0, [
    'Contributions:',
    '- General mathematical framework (6 props)',
    '- Two-layer descriptor methodology',
    '- Data four-classification theory',
    '- State-conditioned expert routing',
    '',
    'Experiments:',
    '- MLIP (verified) + MedMNIST (verified)',
    '- Synthetic data benchmark',
    '- Planned: FEM simulation case study',
    '',
    'Target: TMLR / AISTATS',
], fs=13, color=GRAY)
tb(s, 0.8, 6.3, 11.7, 0.4, 'Timeline', fs=18, bold=True, color=DARK)
tb(s, 0.8, 6.7, 11.7, 0.5, 'Paper 1 can start now (complete data). Paper 2 needs FEM case + synthetic benchmark before submission.', fs=14, color=GRAY)
notes(s, 'Two papers. Paper 1 SCX-MLIP has complete experimental data and can start writing now. Paper 2 SCX-Theory needs additional FEM case study and synthetic benchmark. Paper 1 targets Nature Comp Sci/npj CM, Paper 2 targets TMLR/AISTATS.')

# ====== SLIDE 13: Closing ======
s = add_slide(prs); add_bg(s, WHITE)
tb(s, 1.0, 1.8, 11.3, 1.2, 'SCX does not replace models. SCX does not replace experts.', fs=34, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tb(s, 1.0, 3.2, 11.3, 0.8, 'SCX judges: in the current state,', fs=24, color=GRAY, align=PP_ALIGN.CENTER)
tb(s, 1.0, 4.0, 11.3, 1.0, 'which data to trust, which expert to trust, which action is worth the cost.', fs=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
hline(s, 4.5, 5.3, 4.3)
tb(s, 1.0, 5.7, 11.3, 0.6, 'The value of data is state-conditioned.', fs=20, color=GRAY, align=PP_ALIGN.CENTER)
tb(s, 1.0, 6.5, 11.3, 0.5, 'Thank you | Questions welcome', fs=18, color=GRAY, align=PP_ALIGN.CENTER)
notes(s, 'Final slide. One sentence summary: SCX does not replace models or experts. It judges, in the current state, which data to trust, which expert to trust, and which action is worth the cost. The core proposition: The value of data is state-conditioned. Thank you.')

# Save
out = 'G:/Xiaogan_Supercomputing_data/SCX/outputs/SCX_Academic_Presentation_2026-06-26.pptx'
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
