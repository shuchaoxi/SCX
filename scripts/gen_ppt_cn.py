#!/usr/bin/env python
"""Generate SCX academic PPTX -- Chinese version."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x1A, 0x56, 0xDB)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x0D, 0x94, 0x8B)
RED_ACCENT = RGBColor(0xDC, 0x26, 0x26)
ORANGE = RGBColor(0xE6, 0x8A, 0x00)
LIGHT_BLUE = RGBColor(0xDB, 0xE9, 0xFA)
HEADER_BG = RGBColor(0x25, 0x69, 0xEB)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
PINK_BG = RGBColor(0xFD, 0xE8, 0xE8)
WARM_BG = RGBColor(0xFE, 0xF3, 0xE6)

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_bg(slide, color=LIGHT_GRAY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def hline(slide, l, t, w):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Pt(3))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()

def tb(slide, l, t, w, h, text, fs=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(fs)
    p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    return tf

def ml(slide, l, t, w, h, lines, fs=16, color=DARK):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(fs); p.font.color.rgb = color; p.space_after = Pt(4)
    return tf

def card(slide, l, t, w, h, fill=WHITE, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def bar(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ====== SLIDE 1: Cover ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 1.5, 1.8, 10.3, 1.2, 'SCX：状态条件数据价值评估框架', fs=40, bold=True, color=DARK)
tb(s, 1.5, 3.0, 10.3, 0.8, 'The value of data is state-conditioned', fs=24, color=BLUE)
tb(s, 1.5, 3.8, 10.3, 0.5, '从机器学习势函数到科学与工程仿真', fs=18, color=GRAY)
tb(s, 1.5, 5.5, 5, 1.0, '汇报人：许少杰\n华中科技大学\n2026 年 6 月', fs=16, color=GRAY)
notes(s, '各位老师好，我今天汇报的题目是 SCX——状态条件数据价值评估框架。出发点很朴素：数据很贵，但数据价值不是样本自身的固定属性，而是取决于状态的。我从机器学习势函数里发现了这个规律，然后把它抽象成了通用框架。今天讲四件事：问题是什么、核心想法、做了什么、下一步怎么做。')

# ====== SLIDE 2: Problem ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, '问题背景：数据很贵，但我们经常不知道哪些数据值得用', fs=26, bold=True, color=DARK)
problems = [
    ('数据采集\n成本高昂', ['DFT 一个结构可能算几十核时', 'FEM 一次仿真跑几小时', '医学图像标注需要专家']),
    ('数据价值\n不明确', ['哪些样本值得继续算？', '哪些是重复的、可以跳过？', '哪些数据可能是错的？']),
    ('多专家/多模型\n难以比较', ['不存在全局最优的模型', '专家可靠性随区域变化', '如何把样本路由到合适的专家？']),
    ('传统方法\n有局限', ['LOO/Shapley 给每个样本固定分数', '不考虑状态上下文', '分不清噪声和困难样本']),
]
for i, (title, items) in enumerate(problems):
    x = 0.5 + i * 3.1
    card(s, x, 1.5, 2.9, 3.8, fill=LIGHT_GRAY)
    tb(s, x + 0.2, 1.6, 2.5, 0.7, title, fs=15, bold=True, color=DARK)
    ml(s, x + 0.2, 2.4, 2.5, 2.8, items, fs=13, color=GRAY)
tb(s, 0.8, 5.7, 11.7, 1.0, '核心矛盾：我们需要一个框架，能在状态的粒度上判断——这个数据值不值得用、这个专家值不值得信、这个动作值不值得花钱。', fs=17, color=DARK)
notes(s, '先说为什么要做这个。在科学计算里数据很贵，但我们经常不知道花大价钱算出来的数据到底哪些真有用、哪些是重复劳动、哪些甚至有问题。传统方法给每个样本一个固定分数，但数据价值不是样本的固有属性。同一个样本，在不同模型、不同已有数据覆盖下，价值完全不一样。打个比方：一道题对刷了100道同类题的学生是冗余的，对刚入门的学生是高价值的。SCX 要解决的就是这个矛盾。')

# ====== SLIDE 3: Core Idea ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, '核心思想：数据价值是状态条件的', fs=28, bold=True, color=DARK)
card(s, 0.5, 1.5, 5.8, 2.0, fill=LIGHT_GRAY)
tb(s, 0.8, 1.6, 5.2, 0.4, '传统方法：全局打分', fs=20, bold=True, color=RED_ACCENT)
ml(s, 0.8, 2.1, 5.2, 1.2, ['LOO / Data Shapley / Influence Function', '每个样本 → 一个固定分数', '不考虑：谁在用、已有啥、当前缺啥', '结果：噪声和困难样本分不清'], fs=14, color=GRAY)
card(s, 6.8, 1.5, 5.8, 2.0, fill=LIGHT_BLUE)
tb(s, 7.1, 1.6, 5.2, 0.4, 'SCX：状态条件打分', fs=20, bold=True, color=BLUE)
ml(s, 7.1, 2.1, 5.2, 1.2, ['V(x) = V(x | 状态, 覆盖度, 专家, 成本, 反馈)', '同一个样本 → 不同条件下价值不同', '数据价值由五个因素共同决定', '可以区分：冗余 / 噪声 / 高价值 / 需路由'], fs=14, color=DARK)
card(s, 0.5, 3.9, 12.3, 1.2, fill=WHITE, border=BLUE)
tb(s, 0.8, 4.0, 11.7, 0.4, '核心公式', fs=16, bold=True, color=BLUE)
tb(s, 0.8, 4.4, 11.7, 0.5, 'V(s) = r_bar(s) x rho(s) x L(s) x [1 - D(s)] x max_m SCX_m(s)', fs=22, bold=True, color=DARK)
tb(s, 0.8, 4.8, 11.7, 0.3, 'r_bar=当前误差  rho=状态出现概率  L=可学习性(非噪声)  D=已有覆盖度  SCX=该状态下最佳专家的可靠性', fs=11, color=GRAY)
tb(s, 0.8, 5.5, 11.7, 1.0, '关键洞察：数据价值不是样本的固有属性。同一批数据，在模型 A 眼里是噪声，在模型 B 手里可能是高价值样本。', fs=17, color=DARK)
notes(s, '这一页是核心思想。传统方法给每个样本一个固定分数。SCX 认为数据价值取决于状态——包括当前误差、状态出现频率、标签一致性、已有覆盖度、以及哪个专家在这个状态最可靠。公式里五个因子是乘积关系：任何一个为零，整个价值就为零。比如当前误差已经很低了，那这个状态的数据就不值钱——模型已经学会了。比如噪声很高，也不值钱——再多的坏数据也学不出好东西。这个思想的来源很朴素：做势函数时发现同一个模型在声子附近预测很准、在高温区预测很差——那显然在声子区继续加数据是浪费。')

# ====== SLIDE 4: Four-quadrant ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, '大白话解释：同一个样本，在不同状态下价值完全不同', fs=26, bold=True, color=DARK)
quads = [
    ('冗余数据', '低误差 + 高密度', '→ 压缩 / 跳过', 'Phonon 120 帧: fmax~0.66\n预估可压缩 50-80%', LIGHT_GRAY),
    ('高价值数据', '高误差 + 高密度 + 高一致性', '→ 继续采集', 'EOS 极端应变区\n弹性变形方向', LIGHT_BLUE),
    ('噪声 / 坏标签', '高误差 + 低密度 + 低一致性', '→ 重算 / 降权', 'Thermal 1800K: 49.1% 帧 fmax>5\n14 帧 fmax>10 (极端异常)', PINK_BG),
    ('专家依赖数据', '某专家在此区域显著优于其他', '→ 路由到最合适的专家', 'MLMD stress10 区域\n需路由到高温训练专家', WARM_BG),
]
for i, (title, cond, action, example, bg) in enumerate(quads):
    x = 0.3 + (i % 2) * 6.5
    y = 1.4 + (i // 2) * 2.8
    card(s, x, y, 6.2, 2.5, fill=bg)
    tb(s, x + 0.2, y + 0.1, 5.8, 0.4, title, fs=18, bold=True, color=DARK)
    tb(s, x + 0.2, y + 0.5, 5.8, 0.4, cond, fs=13, color=BLUE)
    tb(s, x + 0.2, y + 0.95, 5.8, 0.3, action, fs=14, bold=True, color=DARK)
    ml(s, x + 0.2, y + 1.3, 5.8, 1.1, example.split('\n'), fs=12, color=GRAY)
notes(s, '这一页用大白话解释四分类。冗余数据：模型已经学得很好，类似样本很多，再加就是浪费钱。AlN 的 phonon 120 帧全部 fmax 在 0.66 左右，SCX 建议可压缩 50-80%。高价值数据：误差还比较高但一致性好——不是噪声，是真的还没学会。噪声数据：误差高、密度低、一致性低——大概率是数据本身有问题。1800K 的热振动 snapshot 有将近一半 fmax 超过 5，14 帧超过 10——物理上 AlN 晶体里原子受力不可能这么大，这就是脚本生成的坏构型，属于数据中毒。专家依赖：某类数据专家 A 处理得很好、专家 B 很差，这时应该路由给 A 而不是搞平均主义。')

# ====== SLIDE 5: Four-classification ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'SCX 把数据分成四类，每类对应不同动作', fs=28, bold=True, color=DARK)
classes = [
    ('Valuable\n高价值', '高误差 + 高密度\n+ 高一致性 + 低覆盖', '→ 采集更多', 'EOS极端应变、弹性变形\n尚未覆盖的构型区域', BLUE),
    ('Redundant\n冗余', '低误差 + 高密度\n+ 高覆盖', '→ 压缩/跳过', 'Phonon小位移 120 帧\n压缩 50-80% 精度不降', ACCENT),
    ('Noisy\n噪声风险', '高误差 + 低密度\n+ 低一致性', '→ 重算/降权/丢弃', 'Thermal 1800K: 49% 异常\nMLMD stress10: 35% 异常', RED_ACCENT),
    ('Expert-dep.\n专家依赖', '某专家在此状态\n显著优于其他', '→ 路由到该专家', 'MLMD 高温区需专项专家\n表面/缺陷需迁移学习', ORANGE),
]
for i, (name, cond, action, example, clr) in enumerate(classes):
    y = 1.4 + i * 1.45
    bar(s, 0.5, y, 0.08, 1.2, clr)
    tb(s, 0.8, y + 0.05, 2.0, 1.1, name, fs=18, bold=True, color=DARK)
    tb(s, 2.8, y + 0.05, 3.2, 1.1, cond, fs=14, color=GRAY)
    tb(s, 6.2, y + 0.05, 2.5, 1.1, action, fs=16, bold=True, color=clr)
    tb(s, 9.0, y + 0.05, 4.0, 1.1, example, fs=13, color=GRAY)
notes(s, '四分类的具体标准和动作建议。高价值：模型还没学会，且数据质量可靠——值得继续投入计算资源。冗余：模型已经学会了，再算类似结构是浪费钱——SCX 在 AlN 上发现 phonon 120 帧可以压缩一半以上。噪声风险：数据本身有问题，强行训练反而不利——SCX 能自动发现这些帧。专家依赖：别搞平均主义，在状态 A 专家 1 最强就把数据给专家 1。这四个分类不是人为划的，是 SCX 根据状态条件自动判定的。')

# ====== SLIDE 6: Why MLIP first ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, '为什么 MLIP 是 SCX 的第一个根据地', fs=28, bold=True, color=DARK)
reasons = [
    ('1. 天然的状态编码器', 'ACE descriptor 把任意原子构型映射到 182 维基函数空间。构型相似 → 基函数接近。不需要额外训练 embedding 网络。物理本身就给了你编码器。'),
    ('2. 可靠的参考标签', 'DFT 虽然贵，但提供了相对可靠的能量/力/应力标签。有 ground truth，才能画 residual map，才能判断专家好不好。这比纯无监督场景多了一个锚点。'),
    ('3. 天然的多专家场景', '同一个训练集可以训出多个势函数——不同架构、不同初始化、不同数据子集。它们在同一个构型区域上表现天然不同——这正是 SCX 需要对比的。'),
]
for i, (title, body) in enumerate(reasons):
    y = 1.5 + i * 1.6
    card(s, 0.5, y, 12.3, 1.4, fill=LIGHT_GRAY if i % 2 == 0 else LIGHT_BLUE)
    tb(s, 0.8, y + 0.1, 11.7, 0.4, title, fs=20, bold=True, color=DARK)
    tb(s, 0.8, y + 0.55, 11.7, 0.7, body, fs=14, color=GRAY)
tb(s, 0.8, 6.5, 11.7, 0.5, '总结：MLIP = 天然编码器(ACE) + 可靠标签(DFT) + 多专家对比 → SCX 的理想起点', fs=17, bold=True, color=BLUE)
notes(s, '为什么 SCX 最早在 MLIP 领域做出来？三个原因。第一，ACE descriptor 本身就是天然编码器——物理给了你 embedding。第二，DFT 提供可靠但昂贵的标签——你有 ground truth 可以校准。第三，MLIP 天然是多专家场景——同一数据可以训不同势函数，它们在相同区域表现天然不同。这三个条件让 MLIP 成为 SCX 最理想的第一个验证场景。')

# ====== SLIDE 7: SCX-MLIP Progress ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'SCX-MLIP：第一个根据地的进展', fs=28, bold=True, color=DARK)
card(s, 0.5, 1.5, 6.0, 5.3, fill=LIGHT_GRAY)
tb(s, 0.8, 1.6, 5.4, 0.4, '势函数侧：Model B + Gauge Fixing', fs=17, bold=True, color=DARK)
ml(s, 0.8, 2.1, 5.4, 4.5, [
    'Single ACE (1196 参数): 力 RMSE 0.045 eV/A',
    '',
    'Model B shared+correction (1378 参数):',
    '  E = sum[c0*Bi + c_Z*Bi + b_Z]',
    '  shared c0 + 元素修正 c_Z',
    '',
    'Gauge Fixing: 后处理投影方法',
    '  先训练，再投影: c_Z -= g, c_0 += g',
    '  Violation = 4.6e-16 (机器零)',
    '  预测完全不变 (数学恒等式)',
    '',
    '关键认识：gauge fix 的真正价值在势函数',
    '合并 (AlN+GaN)，单独 AlN 不带来额外精度',
], fs=12, color=GRAY)
card(s, 6.8, 1.5, 6.0, 5.3, fill=LIGHT_BLUE)
tb(s, 7.1, 1.6, 5.4, 0.4, 'SCX 侧：两层描述符 + 数据防中毒', fs=17, bold=True, color=DARK)
ml(s, 7.1, 2.1, 5.4, 4.5, [
    '一层 (12-dim MLIP encoder):',
    '  50% 帧挤在一个状态, 噪声 F1 = 0.253',
    '',
    '两层 (= 一层 + ErrorDrivenEncoder):',
    '  L1 粗聚类 → L2 误差子空间细化',
    '  噪声 F1: 0.253 → 0.585 (+131%)',
    '  Top-3 状态捕获 94.6% 噪声帧',
    '  Phonon 120 帧 100% 隔离 (低误差)',
    '',
    '数据防中毒 (Data Poisoning Defense):',
    '  无先验知识发现全部 74 个高噪声帧',
    '  训练 fmax vs 测试误差: r = 0.966',
    '  预估去噪后力 RMSE 降 29-48%',
], fs=12, color=GRAY)
bar(s, 0.5, 6.9, 4.0, 0.5, BLUE); tb(s, 0.7, 6.92, 3.6, 0.4, '534 帧 → 可压缩约150帧 + 移除14帧极端噪声', fs=13, bold=True, color=WHITE)
bar(s, 4.8, 6.9, 3.5, 0.5, RED_ACCENT); tb(s, 5.0, 6.92, 3.1, 0.4, 'SCX 100% 命中 REPORT 最差12帧', fs=13, bold=True, color=WHITE)
bar(s, 8.6, 6.9, 4.2, 0.5, ACCENT); tb(s, 8.8, 6.92, 3.8, 0.4, '预估力RMSE: 0.045 → 0.023-0.032 eV/A', fs=13, bold=True, color=WHITE)
notes(s, '第一个根据地的详细进展。左边是势函数工作：Single ACE 基线、Model B shared+correction 架构、以及 gauge fixing。我们发现后处理投影是最优方案——先训练再投影，violation 精确到 10 的负 16 次方，预测完全不变。右边是 SCX 工作：一层方法太粗导致一半帧挤在一个状态，两层方法把噪声 F1 从 0.253 提升到 0.585。最关键的发现是数据防中毒：SCX 不知道 fmax 阈值也不知道 batch 标签，但独立发现了全部 74 个高噪声帧。训练 fmax 和测试误差的 Pearson 相关系数高达 0.966。预估去噪后力 RMSE 降 29-48%。底部三个数字是核心证据。')

# ====== SLIDE 8: MLIP -> SCX-Sim ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, '从 MLIP 到科学仿真：问题结构是一样的', fs=28, bold=True, color=DARK)
tb(s, 0.8, 1.3, 11.7, 0.4, '核心映射：状态 = 仿真参数空间中的区域 | 专家 = 不同保真度的求解器 | 数据价值 = 是否值得在这个参数区域做高保真计算', fs=15, color=BLUE)
mappings = [
    ('MLIP [已验证]', ['状态 = 原子构型空间', '专家 = 不同势函数', '高保真 = DFT 计算', '低保真 = 经验势 / ACE', '噪声 = fmax 异常的帧', '冗余 = 过多相似构型'], True),
    ('FEM / 工程仿真', ['状态 = 参数空间(载荷/边界)', '专家 = 粗/细网格求解器', '高保真 = 精细网格 FEM', '低保真 = 粗网格 / 降阶模型', '噪声 = 数值发散 / 不收敛', '冗余 = 参数扫描过密区'], False),
    ('通用 SCX-Sim', ['状态 = 仿真场景参数组合', '专家 = 各保真度求解器链', '高保真 = 实验 / 高精度仿真', '低保真 = surrogate / 经验公式', '噪声 = 收敛失败 / 数值异常', '冗余 = 参数空间过度采样'], False),
]
for i, (title, items, is_verified) in enumerate(mappings):
    x = 0.5 + i * 4.2
    card(s, x, 2.0, 3.9, 4.0, fill=LIGHT_BLUE if is_verified else LIGHT_GRAY)
    tb(s, x + 0.2, 2.1, 3.5, 0.4, title, fs=15, bold=True, color=DARK)
    ml(s, x + 0.2, 2.6, 3.5, 3.2, items, fs=13, color=DARK if is_verified else GRAY)
tb(s, 0.8, 6.4, 11.7, 0.8, 'MLIP 中发现的规律（噪声集中在极端参数区、冗余在低误差区、专家分歧在特定状态）在 FEM 和工程仿真中同样成立。', fs=15, color=DARK)
notes(s, '从 MLIP 推广到科学仿真的逻辑。仔细看科学和工程仿真里的问题结构，和 MLIP 完全一样。状态：MLIP 里是构型空间，FEM 里是载荷/边界条件组合的参数空间。专家：MLIP 里是不同势函数，仿真里是不同保真度求解器——粗网格便宜但不准，细网格准但贵。噪声：MLIP 里是 fmax 异常，仿真里是数值发散、收敛失败。冗余：MLIP 里是 phonon 小位移的 120 个相似帧，仿真里是参数扫描过密的区域。问题结构完全一样。')

# ====== SLIDE 9: Workflow ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'SCX 通用工作流：五阶段闭环', fs=28, bold=True, color=DARK)
stages = [
    ('1. 编码', '数据 -> 特征向量\n(ACE/CNN/GNN/\n仿真参数向量)', BLUE),
    ('2. 建图', '聚类 -> 状态地图\n(KMeans/谱聚类/\n误差驱动聚类)', HEADER_BG),
    ('3. 估值', '四维评分\n质量 / 冗余 / 噪声\n/ 专家可靠性', ACCENT),
    ('4. 决策', '动作选择\n采集 / 压缩 / 重标\n/ 路由 / 丢弃', ORANGE),
    ('5. 反馈', '真实结果回传\n重校准阈值\n持续更新', PURPLE),
]
for i, (title, desc, clr) in enumerate(stages):
    x = 0.3 + i * 2.6
    card(s, x, 1.6, 2.3, 3.0, fill=WHITE, border=clr)
    bar(s, x, 1.6, 2.3, 0.55, clr)
    tb(s, x + 0.1, 1.63, 2.1, 0.5, title, fs=18, bold=True, color=WHITE)
    ml(s, x + 0.15, 2.3, 2.0, 2.1, desc.split('\n'), fs=12, color=DARK)
    if i < 4:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.35), Inches(2.8), Inches(0.22), Inches(0.4))
        arr.fill.solid(); arr.fill.fore_color.rgb = clr; arr.line.fill.background()
tb(s, 0.8, 5.0, 11.7, 0.4, '反馈闭环：第5步的真实结果（训练loss变化、仿真验证误差）-> 回传校准第3步的阈值 -> 持续改进决策质量', fs=14, color=GRAY)
card(s, 0.5, 5.6, 12.3, 1.7, fill=LIGHT_GRAY)
tb(s, 0.8, 5.7, 11.7, 0.4, '两层描述符（核心方法创新）', fs=16, bold=True, color=DARK)
ml(s, 0.8, 6.1, 5.5, 1.0, ['Layer 1（人工域知识）：手工描述符 → 粗聚类', '  MLIP: coord_num, bond_len, volume...', '  仿真: Re, Ma, 载荷范围, 几何参数...'], fs=12, color=GRAY)
ml(s, 6.8, 6.1, 5.5, 1.0, ['Layer 2（算法发现）：误差驱动特征选择 → 细化', '  互信息筛选误差相关维度', '  仅在误差相关子空间中聚类', '  自动发现 noise/value/redundant 子状态'], fs=12, color=BLUE)
notes(s, 'SCX 通用工作流，五阶段闭环。编码：不管什么类型的数据映射成特征向量。建图：聚类形成状态地图。估值：对每个状态做四维评分。决策：给出动作建议。反馈：真实结果回传校准阈值。底部两层描述符是最大的方法创新：L1 用人的域知识做粗聚类，L2 用误差驱动的特征选择做细化。本质思想是——状态应该由模型在哪里失败来定义，不是由人的直觉来定义。')

# ====== SLIDE 10: Applications ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, 'SCX 的潜在应用方向：同一个问题结构', fs=28, bold=True, color=DARK)
tb(s, 0.8, 1.2, 11.7, 0.4, '核心问题（普适的）：在当前状态 s 下，这个数据 / 这个模型是否可信？', fs=16, color=BLUE)
apps = [
    ('材料 / DFT / MLIP [已验证]', '势函数选择、数据质量筛查、训练集压缩', BLUE),
    ('FEM / 工程仿真', '粗/细网格切换策略、仿真点冗余检测、数值异常筛查', ACCENT),
    ('制造在线监测', '传感器数据质控、异常检测(非统计型)、多模型路由', HEADER_BG),
    ('医学图像质控', '标注质量评估、困难样本 vs 标注错误区分、冗余切片筛选', ORANGE),
    ('机器人轨迹数据', '仿真轨迹质量评估、sim-to-real 场景选择、多控制器路由', PURPLE),
    ('智能驾驶场景数据', 'Corner case 价值评估、仿真场景冗余检测、多传感器融合路由', RED_ACCENT),
]
for i, (name, detail, clr) in enumerate(apps):
    y = 1.8 + i * 0.82
    bar(s, 0.5, y, 0.07, 0.62, clr)
    tb(s, 0.8, y + 0.02, 3.5, 0.55, name, fs=15, bold=True, color=DARK)
    tb(s, 4.5, y + 0.02, 8.3, 0.55, detail, fs=14, color=GRAY)
tb(s, 0.8, 6.9, 11.7, 0.5, '坦诚声明：除 MLIP 外，其他方向尚未验证。列出是为了说明问题结构相同，不是宣称已经解决。', fs=13, bold=False, color=RED_ACCENT)
notes(s, '潜在应用方向。我必须坦诚——除了 MLIP 之外，其他方向都还没有实际验证。列出它们是想说明：这些看似不相关的领域，底层的数据价值判断逻辑是相同的。都是问同一个问题：在当前这个状态下，这个数据可不可信？这个模型可不可信？FEM 里粗网格和细网格的切换策略、制造监测里传感器的异常检测、医学图像里标注质量评估、自动驾驶里 corner case 价值判断——问题结构一样。但再次强调：我们扎实做好的只有 MLIP 这一个。')

# ====== SLIDE 11: Paper Value ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, '为什么 SCX 值得写成论文', fs=28, bold=True, color=DARK)
pillars = [
    ('理论新意', '数据价值\n是状态条件的', ['不是全局标量，是条件概率', '6 个命题覆盖压缩/路由/两层发现', 'Prop 1 (全局排名不充分性) 已证明', '重新定义了数据价值的数学对象'], BLUE),
    ('应用验证', 'MLIP 场景\n已有量化证据', ['AlN v3: 噪声F1=0.585, Top-3 捕获94.6%', '数据防中毒: 100%命中, r=0.966', 'MedMNIST: 30%压缩精度+6%', 'MedMNIST: SCX路由优于均匀集成'], ACCENT),
    ('通用性', '不限于任何\n特定ML范式', ['不依赖特定模型架构', '不依赖梯度/损失函数形式', '抽象基类 + YAML 域配置', '新领域 3 步接入'], ORANGE),
    ('实用性', '直接回答\n值不值得花钱', ['数据采集: 哪些区域值得继续投入', '数据清洗: 自动发现噪声/异常', '专家部署: 路由到最可靠模型', '成本优化: 压缩冗余, 聚焦高价值'], PURPLE),
]
for i, (title, subtitle, points, clr) in enumerate(pillars):
    x = 0.3 + i * 3.25
    card(s, x, 1.5, 3.0, 5.2, fill=LIGHT_GRAY)
    bar(s, x, 1.5, 3.0, 0.8, clr)
    tb(s, x + 0.15, 1.52, 2.7, 0.4, title, fs=18, bold=True, color=WHITE)
    tb(s, x + 0.15, 1.92, 2.7, 0.3, subtitle, fs=11, color=WHITE)
    ml(s, x + 0.15, 2.5, 2.7, 4.0, points, fs=12, color=DARK if i == 0 else GRAY)
notes(s, '四个角度论证 SCX 的论文价值。理论新意：重新定义了数据价值这个概念——不是全局标量，是状态条件概率。应用验证：AlN v3 和 MedMNIST 上都有量化证据。通用性：不绑定任何特定模型架构，抽象基类加 YAML 配置，新领域三步接入。实用性：直接回答很实际的问题——这个数据值不值得花钱。')

# ====== SLIDE 12: Paper Roadmap ======
s = add_slide(prs); add_bg(s, WHITE); hline(s, 0, 0, 13.333)
tb(s, 0.8, 0.5, 11.7, 0.7, '计划中的两篇论文', fs=28, bold=True, color=DARK)
card(s, 0.5, 1.5, 5.8, 4.5, fill=LIGHT_BLUE, border=BLUE)
tb(s, 0.8, 1.6, 5.2, 0.5, '论文一：SCX-MLIP', fs=22, bold=True, color=BLUE)
tb(s, 0.8, 2.1, 5.2, 0.4, '专家势函数合并与状态条件数据价值', fs=15, color=DARK)
ml(s, 0.8, 2.7, 5.2, 3.0, [
    '核心贡献：',
    '- Model B shared+correction 架构',
    '- Post-hoc gauge fixing 方法',
    '- SCX 两层噪声/冗余检测 (AlN v3)',
    '- 数据防中毒 (data poisoning defense)',
    '',
    '实验基础：',
    '- AlN v3: 534 DFT帧, 689 清洁标签',
    '- Single ACE vs Model B 完整对比',
    '- SCX 噪声检测 100% 命中验证',
    '',
    '目标期刊：Nature Comp Sci / npj CM',
], fs=13, color=GRAY)
card(s, 6.8, 1.5, 5.8, 4.5, fill=LIGHT_GRAY, border=ACCENT)
tb(s, 7.1, 1.6, 5.2, 0.5, '论文二：SCX-Theory', fs=22, bold=True, color=ACCENT)
tb(s, 7.1, 2.1, 5.2, 0.4, 'The value of data is state-conditioned', fs=15, color=DARK)
ml(s, 7.1, 2.7, 5.2, 3.0, [
    '核心贡献：',
    '- 通用数学框架 (6 个命题)',
    '- 两层描述符方法论',
    '- 数据四分类理论',
    '- 状态条件专家路由',
    '',
    '实验基础：',
    '- MLIP (已验证) + MedMNIST (已验证)',
    '- 合成数据 benchmark',
    '- 计划: FEM 仿真案例',
    '',
    '目标期刊：TMLR / AISTATS',
], fs=13, color=GRAY)
tb(s, 0.8, 6.3, 11.7, 0.4, '写作节奏', fs=18, bold=True, color=DARK)
tb(s, 0.8, 6.7, 11.7, 0.5, '论文一可先行（已有完整实验数据）。论文二需要额外完成 FEM 仿真案例 + 合成 benchmark 后投稿。', fs=14, color=GRAY)
notes(s, '计划写两篇论文。论文一 SCX-MLIP 聚焦势函数合并和数据质量——已有完整实验数据，现在就可以开始写。论文二 SCX-Theory 推广为通用框架——需要额外的 FEM 案例和合成 benchmark。时间线上，论文一先行。')

# ====== SLIDE 13: Closing ======
s = add_slide(prs); add_bg(s, WHITE)
tb(s, 1.0, 1.8, 11.3, 1.2, 'SCX 不替代模型，也不替代专家。', fs=36, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tb(s, 1.0, 3.2, 11.3, 0.8, 'SCX 判断：在当前状态下，', fs=24, color=GRAY, align=PP_ALIGN.CENTER)
tb(s, 1.0, 4.0, 11.3, 1.0, '哪些数据值得信、哪些专家值得信、哪些动作值得花钱。', fs=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
hline(s, 4.5, 5.3, 4.3)
tb(s, 1.0, 5.7, 11.3, 0.6, 'The value of data is state-conditioned.', fs=20, color=GRAY, align=PP_ALIGN.CENTER)
tb(s, 1.0, 6.5, 11.3, 0.5, '谢谢各位老师 | 欢迎提问', fs=18, color=GRAY, align=PP_ALIGN.CENTER)
notes(s, '最后一页，一句话总结。SCX 不替代任何模型，也不替代任何专家。它做的事情很朴素——判断在当前这个状态下，哪些数据值得你信、哪些专家值得你信、哪些动作值得你花钱。核心命题只有一句话：The value of data is state-conditioned。数据价值是状态条件的。谢谢各位老师，欢迎提问。')

# Save
out = 'G:/Xiaogan_Supercomputing_data/SCX/outputs/SCX_Academic_Presentation_CN_2026-06-26.pptx'
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
